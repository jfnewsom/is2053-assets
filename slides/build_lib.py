"""
build_lib.py - Shared deck assembly library for IS2053 slide decks.

Per-deck build scripts (e.g., slides/m1_u1/build_m1_u1.py) import DeckBuilder
and the layout-specific add_* methods, then make a sequence of calls with the
deck's content inline as string literals. DeckBuilder handles the template
unpack, layout duplication, content population, PNG render-and-insert, speaker
notes, and final packaging.

PATTERN
-------
    from build_lib import DeckBuilder, format_concept_notes, format_demo_notes

    deck = DeckBuilder(template='/path/to/IS2053_Template_v2.pptx',
                       png_out='/home/claude/m1_u1_pngs',
                       work='/tmp/build_m1_u1')

    deck.add_title_slide(
        main_title='The Chapter 2 Toolkit',
        subtitle='IS2053 Programming I  Module 1  Unit 1',
        notes='Cold open. Brief greeting then advance.',
    )

    deck.add_narrative_slide(
        title='The Chapter 2 Toolkit',
        section_label='THE BIG IDEA',
        bold_oneliner='From a closed box to a working program.',
        body='A program with no inputs and no outputs is a closed box...',
        notes=format_concept_notes(video_script='...', think_about=[...]),
    )

    deck.add_concept_slide(
        title='Why We Need Variables',
        bullets=['Bullet 1.', 'Bullet 2.', 'Bullet 3.', 'Bullet 4.', 'Bullet 5.'],
        notes=format_concept_notes(...),
    )

    deck.add_demo_slide(
        title='Demo: Asking for a Name',
        code='print("hello")',
        png='slide06a_input_demo.png',
        notes=format_demo_notes(code='print("hello")', instructor='Type slowly.'),
    )

    # ... more slides ...

    deck.save('/path/to/M1-U1.pptx')


LAYOUT TO TEMPLATE SOURCE MAPPING (v2 template, 35 source layouts at the back)
-------------------------------------------------------------------------------
    title slide      -> template slide 1
    narrative intro  -> template slide 33
    concept          -> template slide 2
    demo             -> template slide 3
    output           -> template slide 4
    two-column       -> template slide 23
    overview         -> template slide 20

When the v2 template gains new layouts, add the layout source index to LAYOUT
and write a new add_*_slide method.
"""
import os
import re
import shutil
import subprocess
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

# ----- Paths -----
LIB_DIR = Path(__file__).resolve().parent
REPO_ROOT = LIB_DIR.parent
TOOLS_DIR = REPO_ROOT / 'tools'
DEFAULT_TEMPLATE = REPO_ROOT / 'ppt' / 'IS2053_Template_v2.pptx'

# Make the renderer importable
sys.path.insert(0, str(TOOLS_DIR))
from render_slides import render_code, render_output  # noqa: E402

# pptx skill scripts
PPTX_SCRIPTS = Path('/mnt/skills/public/pptx/scripts')
UNPACK = PPTX_SCRIPTS / 'office' / 'unpack.py'
PACK = PPTX_SCRIPTS / 'office' / 'pack.py'
CLEAN = PPTX_SCRIPTS / 'clean.py'
ADD_SLIDE = PPTX_SCRIPTS / 'add_slide.py'

# ----- Template layout source indices (1-based template slide numbers) -----
LAYOUT = {
    'title':     1,
    'narrative': 33,
    'concept':   2,
    'demo':      3,
    'output':    4,
    'two_col':   23,
    'overview':  20,
}

# ----- Placeholder text in each template layout (used for surgical replacement) -----
# These strings live in the v2 template; if the template changes, update here.
PH = {
    'title': {
        'main':     'Slide Deck Title Goes Here',
        'subtitle': 'IS2053 Programming I  \u2022  Module N  \u2022  Unit N',
    },
    'narrative': {
        'title':         'Narrative Intro',
        'section_label': 'THE ENGAGEMENT',
        'bold_oneliner': 'Build a real game for a real client.',
        'body':          (
            "You're a Junior Developer at Bat City Collective, an Austin indie studio. "
            "Over five modules you'll ship a retro text adventure called Deep in the Heart: "
            "A Lone Star Journey for your client, All My Eggses Live in Texas. Every Python "
            "concept you learn serves the game."
        ),
    },
    'concept': {
        'title':   'Slide Title \u2014 Concept',
        'bullets': [
            'First key point about the concept being introduced',
            'Second point \u2014 what this replaces or improves',
            'Third point \u2014 connect to the Deep in the Heart game',
            'Fourth point \u2014 Think About This: a question for the student',
            'Fifth point \u2014 the takeaway before the demo',
        ],
    },
    'demo': {
        'title': 'Demo: Topic Name',
    },
    'output': {
        'title': 'Output: Topic Name',
    },
    'two_col': {
        'title':         'Two-Column Comparison',
        'left_label':    'WHAT THIS IS FOR',
        'left_bullets':  [
            'Parallel point one \u2014 keep parallel with the right side',
            'Parallel point two \u2014 short and direct',
            'Parallel point three \u2014 about one line each',
            'Parallel point four \u2014 four to six total',
        ],
        'right_label':   "WHAT THIS ISN'T",
        'right_bullets': [
            'Counter point one \u2014 what the left side rules out',
            'Counter point two \u2014 the misconception to kill',
            'Counter point three \u2014 what gets confused with it',
            'Counter point four \u2014 the line that defines it',
        ],
    },
    'overview': {
        'title':            'Overview',
        'section1_label':   'WELCOME',
        'section1_body':    (
            "This semester you'll learn Python by building a complete video game for a "
            "fictional client. No prior programming experience required \u2014 just curiosity "
            "and a willingness to work through problems one step at a time."
        ),
        'section2_body':    (
            "Our textbook is Starting Out with Python, 6th Edition by Tony Gaddis. We "
            "cover Chapters 2\u201311 across five modules, each building on the last. Every "
            "lab adds to the same codebase \u2014 by the end of the semester you'll have a "
            "working game."
        ),
        'section3_label':   'WHAT YOU WILL BUILD',
        'section3_body':    (
            "Deep in the Heart: A Lone Star Journey \u2014 a text-based Texas road trip "
            "game built incrementally across five modules using Python. Each module "
            "introduces new concepts that get applied directly to the game, so your code "
            "always does something real."
        ),
    },
}

# ----- PNG placement (matches what the template prescribes on demo/output slides) -----
PNG_X = Inches(0.90)
PNG_Y = Inches(1.90)
PNG_W = Inches(11.53)
PNG_PLACEHOLDER_TEXTS = ('CODE DEMO', 'EXPECTED OUTPUT', '[ PNG GOES HERE ]', 'Place rendered PNG at:')


# =============================================================================
# Helpers
# =============================================================================

def _xml_escape(s):
    return s.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')


def _replace_text(xml, old_plain, new_plain):
    """Replace plaintext inside an <a:t> tag, escaping the new content."""
    old_escaped = _xml_escape(old_plain)
    new_escaped = _xml_escape(new_plain)
    pattern = rf'(<a:t(?:\s[^>]*)?>){re.escape(old_escaped)}(</a:t>)'
    new_xml, count = re.subn(pattern, rf'\g<1>{new_escaped}\g<2>', xml)
    if count == 0:
        raise ValueError(f'Placeholder not found in slide XML: {old_plain!r}')
    return new_xml


# =============================================================================
# Notes formatters
# =============================================================================

def format_concept_notes(video_script, key_terms=None, think_about=None, source_url=None):
    """Compose a notes string from the standard concept-slide pieces.

    key_terms: list of (term, definition) tuples or None.
    think_about: list of question strings or None.
    source_url: full URL string or None.
    """
    parts = ['VIDEO SCRIPT', '', video_script.strip()]
    if key_terms:
        parts += ['', 'KEY TERMS']
        parts += [f'{t}: {d}' for t, d in key_terms]
    if think_about:
        parts += ['', 'THINK ABOUT THIS']
        parts += [f'{i+1}. {q}' for i, q in enumerate(think_about)]
    if source_url:
        parts += ['', f'SOURCE: {source_url}']
    return '\n'.join(parts)


def format_demo_notes(code, instructor_notes):
    """Compose demo-slide notes: the code (re-shown in notes) plus directing notes."""
    return f'DEMO CODE (type this live):\n\n{code.strip()}\n\n\nINSTRUCTOR NOTES\n\n{instructor_notes.strip()}'


def format_output_notes(output_text, instructor_notes):
    """Compose output-slide notes."""
    return f'EXPECTED OUTPUT\n\n{output_text.strip()}\n\n\nINSTRUCTOR NOTES\n\n{instructor_notes.strip()}'


def format_title_notes(deck_id, deck_title, opening_line=None):
    """Compose title-slide notes: a brief cold-open guide."""
    default_opening = (
        f'Welcome to {deck_id}. Brief greeting, name the unit, then advance.'
    )
    return (
        f'TITLE SLIDE \u2014 {deck_id}: {deck_title}\n\n'
        f'This is the cold open. Brief greeting, name the unit, then immediately advance.\n\n'
        f'Suggested opening line:\n'
        f'{opening_line or default_opening}'
    )


# =============================================================================
# DeckBuilder
# =============================================================================

class DeckBuilder:
    """Assembles a slide deck from per-layout add_*_slide calls.

    Lifecycle:
        deck = DeckBuilder(template=..., png_out=..., work=...)
        deck.add_<layout>_slide(...)   # repeat for each slide
        deck.save(output_path)
    """

    def __init__(self, template=None, png_out=None, work=None):
        self.template = Path(template or DEFAULT_TEMPLATE)
        if not self.template.exists():
            raise FileNotFoundError(f'Template not found: {self.template}')

        self.png_out = Path(png_out or '/home/claude/build_pngs')
        self.png_out.mkdir(parents=True, exist_ok=True)

        self.work = Path(work or '/tmp/build_deck')
        if self.work.exists():
            shutil.rmtree(self.work)
        self.work.mkdir(parents=True)

        # Unpack the template
        self._run([sys.executable, str(UNPACK), str(self.template), str(self.work / 'unpacked')])
        self.unpacked = self.work / 'unpacked'
        self.slides_dir = self.unpacked / 'ppt' / 'slides'

        # Track new slides and pending PNG inserts / notes
        # Each slide entry: {'xml_file': 'slide36.xml', 'notes': '<text>', 'png': <png path or None>}
        self.slides = []

    # ----- subprocess helper -----
    def _run(self, cmd):
        result = subprocess.run(cmd, capture_output=True, text=True)
        if result.returncode != 0:
            raise RuntimeError(f'Command failed: {cmd}\n{result.stderr}')
        return result.stdout

    # ----- internal: duplicate a template layout and return the new slide filename -----
    def _duplicate(self, layout_key):
        source_slide = f'slide{LAYOUT[layout_key]}.xml'
        out = self._run([sys.executable, str(ADD_SLIDE), str(self.unpacked), source_slide])
        m = re.search(r'Created (slide\d+\.xml)', out)
        if not m:
            raise RuntimeError(f'Could not parse add_slide output:\n{out}')
        return m.group(1)

    def _read(self, slide_filename):
        return (self.slides_dir / slide_filename).read_text()

    def _write(self, slide_filename, xml):
        (self.slides_dir / slide_filename).write_text(xml)

    # =========================================================================
    # Slide builders (one per layout)
    # =========================================================================

    def add_title_slide(self, main_title, subtitle, notes=''):
        """Title slide. main_title goes in the big centered slot; subtitle in the small slot."""
        fn = self._duplicate('title')
        xml = self._read(fn)
        xml = _replace_text(xml, PH['title']['main'], main_title)
        xml = _replace_text(xml, PH['title']['subtitle'], subtitle)
        self._write(fn, xml)
        self.slides.append({'xml_file': fn, 'notes': notes, 'png': None, 'is_png_slide': False})

    def add_narrative_slide(self, title, section_label, bold_oneliner, body, notes=''):
        """Narrative intro. section_label is the small caps tag (e.g., THE BIG IDEA)."""
        fn = self._duplicate('narrative')
        xml = self._read(fn)
        xml = _replace_text(xml, PH['narrative']['title'], title)
        xml = _replace_text(xml, PH['narrative']['section_label'], section_label)
        xml = _replace_text(xml, PH['narrative']['bold_oneliner'], bold_oneliner)
        xml = _replace_text(xml, PH['narrative']['body'], body)
        self._write(fn, xml)
        self.slides.append({'xml_file': fn, 'notes': notes, 'png': None, 'is_png_slide': False})

    def add_concept_slide(self, title, bullets, notes=''):
        """Concept slide: title + exactly 5 bullets. Pad with empty strings if you have fewer."""
        if len(bullets) > 5:
            raise ValueError(f'Concept slide accepts up to 5 bullets, got {len(bullets)}.')
        bullets = list(bullets) + [''] * (5 - len(bullets))

        fn = self._duplicate('concept')
        xml = self._read(fn)
        xml = _replace_text(xml, PH['concept']['title'], title)
        for ph, bullet in zip(PH['concept']['bullets'], bullets):
            xml = _replace_text(xml, ph, bullet)
        self._write(fn, xml)
        self.slides.append({'xml_file': fn, 'notes': notes, 'png': None, 'is_png_slide': False})

    def add_demo_slide(self, title, code, png, notes=''):
        """Demo (code) slide. Renders the PNG from `code` into self.png_out/<png>."""
        fn = self._duplicate('demo')
        xml = self._read(fn)
        xml = _replace_text(xml, PH['demo']['title'], title)
        self._write(fn, xml)

        render_code(code=code, label='', filename=png, out_dir=str(self.png_out))
        self.slides.append({
            'xml_file': fn, 'notes': notes,
            'png': str(self.png_out / png), 'is_png_slide': True,
        })

    def add_output_slide(self, title, output_text, png, notes=''):
        """Output slide. Renders the PNG from `output_text` into self.png_out/<png>."""
        fn = self._duplicate('output')
        xml = self._read(fn)
        xml = _replace_text(xml, PH['output']['title'], title)
        self._write(fn, xml)

        render_output(text=output_text, label='', filename=png, out_dir=str(self.png_out))
        self.slides.append({
            'xml_file': fn, 'notes': notes,
            'png': str(self.png_out / png), 'is_png_slide': True,
        })

    def add_two_column_slide(self, title, left_label, left_bullets, right_label, right_bullets, notes=''):
        """Two-column comparison. Each side accepts up to 4 bullets."""
        if len(left_bullets) > 4 or len(right_bullets) > 4:
            raise ValueError('Two-column slide accepts up to 4 bullets per side.')
        left_bullets = list(left_bullets) + [''] * (4 - len(left_bullets))
        right_bullets = list(right_bullets) + [''] * (4 - len(right_bullets))

        fn = self._duplicate('two_col')
        xml = self._read(fn)
        xml = _replace_text(xml, PH['two_col']['title'], title)
        xml = _replace_text(xml, PH['two_col']['left_label'], left_label)
        for ph, b in zip(PH['two_col']['left_bullets'], left_bullets):
            xml = _replace_text(xml, ph, b)
        xml = _replace_text(xml, PH['two_col']['right_label'], right_label)
        for ph, b in zip(PH['two_col']['right_bullets'], right_bullets):
            xml = _replace_text(xml, ph, b)
        self._write(fn, xml)
        self.slides.append({'xml_file': fn, 'notes': notes, 'png': None, 'is_png_slide': False})

    def add_overview_slide(self, title, section1_label, section1_body, section2_body,
                           section3_label, section3_body, notes=''):
        """Overview / recap slide. Three section bodies, two of them with section labels."""
        fn = self._duplicate('overview')
        xml = self._read(fn)
        xml = _replace_text(xml, PH['overview']['title'], title)
        xml = _replace_text(xml, PH['overview']['section1_label'], section1_label)
        xml = _replace_text(xml, PH['overview']['section1_body'], section1_body)
        xml = _replace_text(xml, PH['overview']['section2_body'], section2_body)
        xml = _replace_text(xml, PH['overview']['section3_label'], section3_label)
        xml = _replace_text(xml, PH['overview']['section3_body'], section3_body)
        self._write(fn, xml)
        self.slides.append({'xml_file': fn, 'notes': notes, 'png': None, 'is_png_slide': False})

    # =========================================================================
    # save: finalize sldIdLst, clean, pack, then inject PNGs + notes
    # =========================================================================

    def save(self, output_path):
        """Update sldIdLst, clean, pack, then add PNGs and speaker notes."""
        # 1) Update sldIdLst to contain only our new slides in order
        pres_xml_path = self.unpacked / 'ppt' / 'presentation.xml'
        pres_xml = pres_xml_path.read_text()

        # The new slides get rels starting from the highest existing rId.
        # add_slide.py appends rels in order, so we need to find each slide's rId.
        # Easiest: parse presentation.xml.rels and map slide{N}.xml -> rId.
        rels_xml = (self.unpacked / 'ppt' / '_rels' / 'presentation.xml.rels').read_text()
        slide_to_rid = {}
        for m in re.finditer(r'Id="(rId\d+)"\s+Type="[^"]*slide"\s+Target="slides/(slide\d+\.xml)"', rels_xml):
            slide_to_rid[m.group(2)] = m.group(1)

        entries = []
        for i, slide in enumerate(self.slides, start=400):
            rid = slide_to_rid.get(slide['xml_file'])
            if rid is None:
                raise RuntimeError(f"Could not find rId for {slide['xml_file']}")
            entries.append(f'    <p:sldId id="{i}" r:id="{rid}"/>')
        new_sldidlst = '<p:sldIdLst>\n' + '\n'.join(entries) + '\n  </p:sldIdLst>'

        pres_xml = re.sub(r'<p:sldIdLst>.*?</p:sldIdLst>', new_sldidlst, pres_xml, flags=re.DOTALL)
        pres_xml_path.write_text(pres_xml)

        # 2) Clean + pack to an intermediate file
        self._run([sys.executable, str(CLEAN), str(self.unpacked)])

        intermediate = self.work / '_intermediate.pptx'
        self._run([sys.executable, str(PACK), str(self.unpacked), str(intermediate),
                   '--original', str(self.template)])

        # 3) Open with python-pptx; add PNGs and notes
        prs = Presentation(str(intermediate))
        if len(prs.slides) != len(self.slides):
            raise RuntimeError(
                f'Slide count mismatch: packed deck has {len(prs.slides)}, expected {len(self.slides)}'
            )

        for i, (slide, spec) in enumerate(zip(prs.slides, self.slides)):
            # Insert PNG on demo/output slides
            if spec['is_png_slide']:
                # Remove placeholder textboxes that would peek through the PNG's rounded corners
                shapes_to_remove = []
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    if any(ph in shape.text_frame.text for ph in PNG_PLACEHOLDER_TEXTS):
                        shapes_to_remove.append(shape)
                for shape in shapes_to_remove:
                    shape._element.getparent().remove(shape._element)

                slide.shapes.add_picture(spec['png'], PNG_X, PNG_Y, width=PNG_W)

            # Speaker notes
            if spec['notes']:
                slide.notes_slide.notes_text_frame.text = spec['notes']

        # 4) Save final
        Path(output_path).parent.mkdir(parents=True, exist_ok=True)
        prs.save(output_path)
        print(f'Saved: {output_path}')
        print(f'  {len(self.slides)} slides, {sum(1 for s in self.slides if s["png"]) } PNGs, '
              f'{sum(1 for s in self.slides if s["notes"]) } notes pages.')
